"""
AI-Powered Proposal Generator — Protiviti DPDPA Privacy Proposal

Changes per run:
  Slide 1  – Company name
  Slide 4  – AI rewrites company description (detailed 3-sentence paragraph) +
             scope paragraph + 7 scope bullets (word-limited to prevent overflow)
             RIGHT-side bullets fixed to circle • matching left side
  Slide 11 – AI rewrites operating model paragraph
  Slide 12, 14, 19 – Company name auto-replaced
  Slide 17 – AI rewrites all 6 Data Lifecycle sections
  Slide 5  – NOT touched (human fills manually)
"""

import io, json, re
from copy import deepcopy
from lxml import etree
import streamlit as st
import streamlit.components.v1 as components
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from groq import Groq

# ─────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────
st.set_page_config(page_title="AI Proposal Generator", page_icon="📊", layout="wide")

# ── Hide "Manage app" bubble — all themes, all devices ──
st.markdown("""
<style>
/* Manage app bar — primary selector */
[data-testid="stStatusWidget"]              { display: none !important; }
[data-testid="stStatusWidget"] *            { display: none !important; }

/* Toolbar top-right (Fork, GitHub) */
[data-testid="stToolbarActions"]            { display: none !important; }
[data-testid="stMainMenuPopover"]           { display: none !important; }

/* Footer */
footer                                      { display: none !important; }
#MainMenu                                   { display: none !important; }
</style>
""", unsafe_allow_html=True)

st.title("📊 AI Proposal Generator")
st.markdown(
    "Upload the **template PPTX**, enter the target company details — "
    "the app personalises every relevant slide while keeping design, fonts and layout intact."
)

# ─────────────────────────────────────────────────────────────
# GROQ API KEY — read from Streamlit secrets first, fall back to manual input
# In Streamlit Cloud: set GROQ_API_KEY in App Settings → Secrets
# Locally: add to .streamlit/secrets.toml as: GROQ_API_KEY = "gsk_..."
# ─────────────────────────────────────────────────────────────
_secret_key = st.secrets.get("GROQ_API_KEY", "") if hasattr(st, "secrets") else ""

# ─────────────────────────────────────────────────────────────
# SIDEBAR
# ─────────────────────────────────────────────────────────────
with st.sidebar:
    st.header("⚙️ Configuration")
    if _secret_key:
        groq_api_key = _secret_key
        st.success("🔑 Groq API key loaded from Secrets", icon="✅")
    else:
        groq_api_key = st.text_input(
            "Groq API Key",
            type="password",
            placeholder="gsk_...",
            help="Or set GROQ_API_KEY in Streamlit Secrets to avoid entering it every time.",
        )
    st.markdown("---")
    st.header("🏢 Company Details")
    company_name    = st.text_input("Full Company Name",         placeholder="SGD Pharma India Pvt. Ltd.")
    company_short   = st.text_input("Abbreviation / Short Name", placeholder="SGD")
    company_website = st.text_input("Website URL",               placeholder="https://www.sgd-pharma.com")
    st.markdown("---")
    st.header("📁 Template")
    uploaded_ppt = st.file_uploader("Upload Template PPTX", type=["pptx"])
    st.markdown("---")
    generate_btn = st.button("🚀 Generate Proposal", use_container_width=True, type="primary")

# ─────────────────────────────────────────────────────────────
# NAMESPACE SHORTHAND
# ─────────────────────────────────────────────────────────────
ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ─────────────────────────────────────────────────────────────
# SCRAPE WEBSITE
# ─────────────────────────────────────────────────────────────
def scrape_website(url: str, max_chars: int = 6000) -> str:
    try:
        r = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=12)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return re.sub(r"\s+", " ", soup.get_text(" ", strip=True))[:max_chars]
    except Exception as e:
        return f"[Scrape failed: {e}]"

# ─────────────────────────────────────────────────────────────
# GROQ CALL
# ─────────────────────────────────────────────────────────────
GROQ_MODEL = "llama-3.3-70b-versatile"

def groq_call(client: Groq, prompt: str, max_tokens: int = 1000) -> str:
    r = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.25,
        max_tokens=max_tokens,
    )
    return r.choices[0].message.content.strip()

# ─────────────────────────────────────────────────────────────
# AI PROMPTS
# ─────────────────────────────────────────────────────────────

# ── Pre-step 1: Extract rich company profile ─────────────────
P_COMPANY_PROFILE = """
You are a business analyst. Read the company website content and extract a structured profile.
Be SPECIFIC — use actual names from the website, not generic placeholders.

Company: {company_name} ({company_short})
Website content: {website_text}

Return a JSON object with EXACTLY these keys:
{{
  "industry": "Primary industry sector e.g. Pharmaceutical Manufacturing, IT Services, BPO, FMCG, Banking",
  "business_model": "One sentence: B2B/B2C/B2B2C split, who they sell to, how they reach clients",
  "service_lines": "Comma-separated 4-6 core service lines or product categories from the website",
  "key_sectors": "Comma-separated 3-5 industry verticals or client sectors served",
  "key_functions": "Comma-separated 5-7 internal business functions e.g. R&D, Manufacturing, Quality, HR, Finance, IT, Regulatory Affairs",
  "key_systems": "Comma-separated 4-6 actual technology systems e.g. ERP, LIMS, CRM, cloud platforms",
  "data_types": "Comma-separated 4-6 main personal data categories e.g. employee data, patient data, client data",
  "partner_types": "Comma-separated 3-5 external party types e.g. distributors, CROs, regulatory bodies, technology vendors",
  "geographic_footprint": "One phrase e.g. pan-India with global exports, nationwide with GCCs",
  "channel_description": "How they primarily operate: direct sales / channel partners / online / B2B contracts etc."
}}

Return ONLY valid JSON. No markdown fences.
"""

# ── Pre-step 2: Business model (kept for compatibility) ──────
P_BIZ_MODEL = """
Analyse the company website content. Answer in ONE sentence (max 20 words):
What is this company's primary business model and channels?

Company: {company_name}
Website content: {website_text}

Return ONLY the single sentence. No labels.
"""

# ── Slide 4: Company description body (TextBox 8 R1) ─────────
# APPROACH: Rewrite ORIGINAL sentence by sentence keeping EXACT grammatical skeleton.
# Only replace: company type, product/service names, channel names, geography.
P_DESC = """
You are rewriting a company description for a professional consulting proposal.
Rewrite the ORIGINAL below for the TARGET COMPANY by SURGICALLY replacing only the
company-specific terms. Keep the EXACT sentence structure, tone, punctuation pattern
and approximately the same word count per sentence.

WHAT TO REPLACE:
- "Indian manufacturer of portable energy and lighting solutions" -> accurate industry/type description
- Product names (batteries, flashlights, lighting) -> actual service lines or products from profile
- "B2B and B2B2C-driven model, serving distributors, retailers, institutional buyers and
  large-scale channel partners through one of India's widest FMCG-style distribution networks" ->
  actual business model and channels from profile
- "brand engagement, after-sales support and product service programs" -> actual B2C touchpoints
- "end-to-end product development, high-volume manufacturing, nationwide distribution and
  lifecycle management" -> actual core delivery capabilities
- "DSIR-approved R&D capabilities, integrated manufacturing facilities" -> actual capabilities/systems
- "portable power and lighting solutions across diverse consumer and commercial segments" ->
  actual value proposition
- "EIIL" -> {company_short}

WHAT TO KEEP EXACTLY (these frames are mandatory):
- Frame 1: "is a leading ... operating through a ... model spanning ... across domestic and
  select international markets."
- Frame 2: "The company follows a predominantly ... model, serving ... through ..., while
  maintaining limited ... interfaces through ..."
- Frame 3: "[Short] enables end-to-end ... through technology-driven quality systems, ...,
  ensuring safe, reliable, compliant and cost-efficient delivery of ... across diverse ...
  segments."

TARGET COMPANY PROFILE:
Name: {company_name} | Short: {company_short}
Industry: {industry}
Business model: {business_model}
Service lines / Products: {service_lines}
Key sectors: {key_sectors}
Key systems: {key_systems}
Geographic footprint: {geographic_footprint}
Channel description: {channel_description}

ORIGINAL (rewrite this — keep skeleton, replace company-specific terms only):
is a leading Indian manufacturer of portable energy and lighting solutions, operating through a
diversified multi-segment model spanning dry-cell batteries, flashlights, consumer lighting,
professional lighting and electrical accessories across domestic and select international markets.
The company follows a predominantly B2B and B2B2C-driven model, serving distributors, retailers,
institutional buyers and large-scale channel partners through one of India's widest FMCG-style
distribution networks, while maintaining limited B2C interfaces through brand engagement,
after-sales support and product service programs. EIIL enables end-to-end product development,
high-volume manufacturing, nationwide distribution and lifecycle management through
technology-driven quality systems, DSIR-approved R&D capabilities, integrated manufacturing
facilities and data-enabled supply-chain operations, ensuring safe, reliable, compliant and
cost-efficient delivery of portable power and lighting solutions across diverse consumer and
commercial segments.

Return ONLY the rewritten body. Start with "is a leading". 118-124 words. No labels, no quotes.
"""

# ── Slide 4: Scope operations phrase (TextBox 3 P0 R3) ───────
P_SCOPE_OPS = """
Complete the sentence below for the target company. Replace [OPS] ONLY — every other word stays.

SENTENCE: "and applicable Rules, calibrated to its people, process and technology landscape
across [OPS]."

[OPS] = 5-10 words describing the company's ACTUAL core operations.
Do NOT use manufacturing/logistics unless the company actually operates in those industries.

Examples:
- BPO/Analytics: "business process outsourcing, analytics and digital services operations"
- Pharma manufacturing: "pharmaceutical manufacturing, R&D and global distribution operations"
- IT services: "enterprise IT services, cloud solutions and managed services operations"
- NBFC: "lending, leasing, digital credit and financial services operations"

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Service lines: {service_lines}
Key sectors: {key_sectors}

Return ONLY the complete sentence starting with "and applicable Rules...". No labels, no quotes.
"""

# ── Slide 4: 7 scope bullets ──────────────────────────────────
P_BULLETS = """
You are rewriting 7 scope bullets for a DPDPA privacy consulting proposal.
APPROACH: Rewrite each ORIGINAL bullet by SURGICALLY replacing only the EIIL/Eveready-specific
operational terms. Keep EVERY privacy/compliance/methodology word IDENTICAL to the original.
Keep the EXACT sentence structure of each bullet.

WHAT TO REPLACE:
- Bullet 1: "EIIL's manufacturing, R&D, supply chain, procurement, commercial, HR, enterprise
  systems and distribution operations" -> replace with actual functions from profile
- Bullet 2: "EIIL's manufacturing, quality, logistics, commercial, enterprise and SaaS platforms"
  -> replace with actual platforms; "distributors, retailers, logistics partners and service
  vendors" -> replace with actual partner types from profile
- Bullet 5: "corporate, manufacturing, R&D, commercial and customer-facing teams" -> replace
  with actual teams from profile
- Bullets 3, 4, 6, 7: Only replace "EIIL" with company short name; keep all else identical

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Service lines: {service_lines}
Key functions: {key_functions}
Key systems: {key_systems}
Partner types: {partner_types}

ORIGINALS (rewrite each — surgical replacement only, keep all privacy words):
1. Conduct an enterprise-wide applicability assessment and privacy gap analysis, covering data discovery, lifecycle mapping, inventories, RoPA and documentation of internal/external data flows across EIIL's manufacturing, R&D, supply chain, procurement, commercial, HR, enterprise systems and distribution operations.
2. Assess privacy, information security and regulatory risks across EIIL's manufacturing, quality, logistics, commercial, enterprise and SaaS platforms, including analytics environments, physical repositories and third-party networks such as distributors, retailers, logistics partners and service vendors.
3. Evaluate governance structures, policies and controls covering lawful purpose, consent (where applicable), retention, erasure, grievance handling, DPR workflows, cross-border transfers and personal data breach processes.
4. Design and operationalize a scalable privacy governance and risk framework, defining roles, accountability, escalation paths and procedures for DPIAs and risk-based reviews of new systems, digital initiatives and operational programs.
5. Support rollout of updated privacy policies, notices and procedures for consent, DPR, retention/deletion, breach response and DPIA processes, tailored for corporate, manufacturing, R&D, commercial and customer-facing teams.
6. Coordinate remediation across key platforms to strengthen consent workflows, DPR handling, third-party data sharing controls, data minimization and privacy-by-design requirements with support from selected tooling partners.
7. Deliver role-based privacy training, define governance KPIs and RACI structures and enable reporting and dashboards to support continuous oversight, audit readiness, regulatory preparedness and executive visibility.

Return exactly 7 lines. No numbering, no bullet symbols, no extra text.
"""

# ── Slide 4: Right-side "How We Will Help" (TextBox 12) ──────
P_S4_RIGHT = """
You are rewriting 6 bullets for the "How We Will Help" section of a DPDPA consulting proposal.
APPROACH: Replace "EIIL" with the company short name. Keep ALL other words IDENTICAL.
These bullets are largely methodology-based so very little changes.

RULES:
- Replace "EIIL" -> {company_short} everywhere
- Keep ALL DPDPA/privacy/governance language WORD FOR WORD
- Match EXACT word count shown in brackets — fixed-size text boxes

TARGET COMPANY: Name: {company_name}, Short: {company_short}
Industry: {industry}

Return JSON with keys "b1","b2","b3","b4","b5","b6". Values = text only. ONLY valid JSON.

ORIGINALS — replace EIIL only, keep everything else word for word:

b1 [EXACTLY 28 words]:
Enable EIIL's transition to sustained compliance with the Digital Personal Data Protection Act by translating regulatory requirements into a risk-calibrated privacy and governance framework aligned to business priorities.

b2 [EXACTLY 22 words]:
Define a risk led, high level compliance roadmap addressing material privacy, data protection and operational gaps across EIIL's personal data processing landscape.

b3 [EXACTLY 16 words]:
Establish prioritized remediation themes, sequencing logic and clear accountability structures to support effective and scalable compliance.

b4 [EXACTLY 25 words]:
Strengthen privacy governance and control architecture by embedding oversight, decision making and process rigor in line with privacy by design and privacy by default principles.

b5 [EXACTLY 25 words]:
Consolidate key observations and the compliance roadmap into formal deliverables to support executive oversight, audit readiness and regulatory preparedness, while enhancing stakeholder trust and transparency.

b6 [EXACTLY 30 words]:
Deliver a risk prioritized remediation roadmap and support governance enablement through a Privacy Steering Committee, defined KPIs, RACI structures and PMO aligned reporting to facilitate coordinated implementation and sustained compliance.
"""

# ── Slide 11: Operating model paragraph ──────────────────────
P_S11 = """
You are rewriting one paragraph for a professional consulting proposal.
The text box is FIXED size — EXACTLY 82 words required.

APPROACH: Rewrite the ORIGINAL by replacing EIIL/Eveready-specific operational terms ONLY.
Keep EVERY other word IDENTICAL to the original. Maintain exact sentence structure.

WHAT TO REPLACE:
- "Eveready Industries India Ltd." -> {company_name}
- "manufacturing, supply-chain, commercial, distribution and corporate operations" ->
  actual operations from service lines and key functions
- "B2B and B2B2C channels" -> actual channels from business model
- "customer service interactions, warranty support and brand engagement programs" ->
  actual B2C data touchpoints (use "client service interactions and digital platform usage"
  if fully B2B)
- "EIIL's" -> {company_short}'s
- "nationwide operational footprint" -> actual geography from profile

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Business model: {business_model}
Service lines: {service_lines}
Key functions: {key_functions}
Geographic footprint: {geographic_footprint}
Channel description: {channel_description}

ORIGINAL (82 words — replace marked terms ONLY, keep all else word for word):
For this engagement, the privacy compliance model will be applied exclusively to the internal
functions, processes and governance structures of Eveready Industries India Ltd., supporting
its manufacturing, supply-chain, commercial, distribution and corporate operations, which
primarily operate through B2B and B2B2C channels, with limited B2C personal data processing
through customer service interactions, warranty support and brand engagement programs. This
approach ensures a focused effort on strengthening EIIL's internal privacy governance and
compliance capabilities, aligned with applicable regulatory requirements, its operating model
and its nationwide operational footprint.

Return ONLY the paragraph. EXACTLY 82 words. No labels, no quotes.
"""

# ── Slide 17: Data Lifecycle ──────────────────────────────────
P_S17 = """
You are rewriting 6 Data Lifecycle paragraphs for a professional consulting proposal.
Each paragraph is in a FIXED-SIZE text box — EXACT word counts are mandatory.

APPROACH: Rewrite each ORIGINAL paragraph by replacing EIIL/Eveready-specific operational
terms with accurate equivalents from the company profile. Keep EVERY other word IDENTICAL.
Maintain the EXACT sentence structure, punctuation and flow of each original.

WHAT TO REPLACE (use the profile below for accurate replacements):
- "Eveready Industries India Ltd. (EIIL)" -> "{company_name} ({company_short})"
- Functions/processes (e.g. "manufacturing processes for batteries, flashlights") ->
  actual functions/service lines from profile
- Systems (e.g. "plant-level manufacturing systems, distribution platforms, logistics systems") ->
  actual systems from profile
- Partner types (e.g. "distributor onboarding") -> actual partner types from profile
- Data types (e.g. "manufacturing, safety") -> actual data categories from profile
- Geography (e.g. "EIIL's nationwide network") -> actual footprint from profile
- "EIIL's" -> "{company_short}'s" everywhere

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Service lines: {service_lines}
Key functions: {key_functions}
Key systems: {key_systems}
Data types: {data_types}
Partner types: {partner_types}
Geographic footprint: {geographic_footprint}

Return JSON keys: "collection","use_processing","storage","sharing","retention","disposal"
Value = paragraph text ONLY (no title, no number). ONLY valid JSON. No markdown fences.

ORIGINALS — surgical replacement only, keep all other words IDENTICAL:

collection [EXACTLY 61 words]:
We will review how Eveready Industries India Ltd. (EIIL) collects personal, operational and regulatory data across functions such as employee onboarding; manufacturing processes for batteries, flashlights and lighting products; distributor onboarding; sales operations; supply-chain coordination; and customer service requirements. This includes data captured through ERP systems, plant-level manufacturing systems, distribution platforms, logistics systems and digital interfaces used across EIIL's nationwide network.

use_processing [EXACTLY 64 words]:
We will assess how collected data is used for manufacturing planning, quality control, inventory management, supply chain coordination, compliance reporting and performance monitoring across EIIL's key segments: batteries, flashlights, consumer lighting, professional lighting and electrical accessories. This includes data integration across systems such as ERP, CRM, distributor management systems and plant-level automation platforms, along with tools supporting R&D operations, workforce management and operational efficiency.

storage [EXACTLY 48 words]:
We will examine secure storage of manufacturing, safety, employee and vendor data across cloud platforms, on-premise servers at EIIL's manufacturing units, validated production systems, backup systems and R&D repositories. Controls for authentication, role-based access, audit trails and compliance with applicable industry and corporate guidelines will also be reviewed.

sharing [EXACTLY 36 words]:
We will evaluate data-sharing practices with distributors, logistics partners, manufacturing vendors, regulatory authorities, retailers and internal teams. This includes reviewing contractual safeguards, supply-chain data-processing requirements, cross-border data transfer practices (where applicable), anonymization procedures and security measures.

retention [EXACTLY 43 words]:
We will review retention policies for manufacturing logs, quality-control reports, product testing data, R&D records, HR and payroll files, vendor documentation, distributor agreements, operational logs and financial documentation. Retention requirements will be assessed against regulatory mandates, audit requirements and internal EIIL governance policies.

disposal [EXACTLY 47 words]:
We will verify secure deletion, destruction and anonymization of records across digital platforms, manufacturing systems, archival repositories, distributor management systems and physical documentation. Disposal workflows will be reviewed for alignment with regulatory expectations and internal EIIL data-governance guidelines to ensure safe and compliant handling of obsolete data.
"""

# ── Slides 12 & 14: Compact operational bullets ──────────────
# Slide 12: 9pt font, 264pt wide box — 32 words (matches original capacity)
# Slide 14: 11pt font, 661pt wide box — 30 words
# PURPOSE: Slide 12 is the CLIENT-FACING SCOPE SUMMARY. Clients read this to understand
# EXACTLY which departments, applications and business functions will be assessed.
# Be SPECIFIC — name real departments, real systems, real business units.
P_S12_BULLETS = """
You are rewriting 2 key scope sentences for a client-facing consulting proposal slide.
This is the most important slide for scope clarity — the client reads these sentences
to understand EXACTLY which departments, systems and business functions will be assessed.

APPROACH: Surgically replace the EIIL-specific operational terms with accurate, SPECIFIC
equivalents for the target company. Keep all privacy/compliance language IDENTICAL.
Be SPECIFIC — name actual departments, business units and systems, not generic terms.

SENTENCE 1 [MAX 32 WORDS — must be complete, specific, end with full stop]:
Original: Conduct an enterprise-wide privacy applicability assessment and gap analysis,
covering data discovery, lifecycle mapping, inventories, RoPA and documentation of
internal/external data flows across EIIL's manufacturing, R&D, supply chain, commercial,
HR, enterprise systems and distribution operations.
→ Replace the operations list with: {company_short}'s actual KEY FUNCTIONS and DEPARTMENTS
  (use 5-6 specific names from the profile, comma-separated, end with "and [last item].")
→ These are the departments the client will know will be in scope.

SENTENCE 2 [MAX 32 WORDS — must be complete, specific, end with full stop]:
Original: Assess privacy, information security and regulatory risks across EIIL's
manufacturing, quality, logistics, commercial, enterprise and SaaS platforms, including
analytics environments, physical repositories and third-party networks such as distributors,
retailers, logistics partners and service vendors.
→ Replace with: {company_short}'s actual KEY SYSTEMS and PLATFORMS (4-5 specific names)
  and actual PARTNER TYPES (2-3 specific types).
→ These tell the client which systems and third parties are in scope.

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Key functions: {key_functions}
Key systems: {key_systems}
Partner types: {partner_types}
Service lines: {service_lines}

Return exactly 2 lines. Line 1 = Sentence 1. Line 2 = Sentence 2.
No numbering, no bullets. Both COMPLETE sentences ending with a full stop.
"""

P_S14_BULLETS = """
Rewrite the 2 sentences below for the TARGET COMPANY.
APPROACH: Surgically replace ONLY the company-specific operational terms.
Keep all privacy/compliance language identical.
Sentences must be COMPLETE and end with a full stop.

SENTENCE 1 [MAX 30 WORDS — complete sentence]:
Original: Conduct an enterprise‑wide privacy applicability assessment and gap analysis,
covering data discovery, lifecycle mapping, inventories, RoPA and documentation of
internal/external data flows across EIIL's manufacturing, R&D, supply chain, commercial,
HR, enterprise systems and distribution operations.
→ Replace "EIIL's manufacturing, R&D, supply chain, commercial, HR, enterprise systems
  and distribution operations" with actual functions — max 5-6 function names.
→ Result: max 30 words, complete sentence.

SENTENCE 2 [MAX 30 WORDS — complete sentence]:
Original: Assess privacy, information security and regulatory risks across EIIL's
manufacturing, quality, logistics, commercial, enterprise and SaaS platforms, including
analytics environments, physical repositories and third-party networks.
→ Replace "EIIL's manufacturing, quality, logistics, commercial, enterprise and SaaS
  platforms" with actual platforms — max 4-5 platform/system names.
→ Result: max 30 words, complete sentence.

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Key functions: {key_functions}
Key systems: {key_systems}
Partner types: {partner_types}

Return exactly 2 lines. Line 1 = Sentence 1. Line 2 = Sentence 2.
No numbering, no bullets. Both sentences must be COMPLETE and end with a full stop.
"""

# ── Slide 19: Privacy Notice sentence (Rectangle 10) ─────────
P_S19_NOTICE = """
Rewrite the sentence below for the TARGET COMPANY. Replace ONLY the list of departments/
functions/platforms after "covering" with accurate equivalents from the company profile.
Keep every other word IDENTICAL.

ORIGINAL:
"Prepare tailored Privacy Notice and Consent Notice for {company_short}'s touchpoints covering
manufacturing, quality, R&D, supply-chain, commercial, HR, enterprise platforms and
cloud/SaaS applications."

RULES:
- Replace only: "manufacturing, quality, R&D, supply-chain, commercial, HR, enterprise
  platforms and cloud/SaaS applications" with actual departments, functions and platforms
  from the company profile — 6 to 8 items, comma-separated, ending with "and [last item]."
- Keep "Prepare tailored Privacy Notice and Consent Notice for {company_short}'s touchpoints
  covering" WORD FOR WORD
- End with a full stop

TARGET COMPANY PROFILE:
Name: {company_name}, Short: {company_short}
Industry: {industry}
Key functions: {key_functions}
Key systems: {key_systems}
Service lines: {service_lines}

Return ONLY the rewritten sentence. No labels, no quotes.
"""

# ─────────────────────────────────────────────────────────────
# PPTX HELPERS
# ─────────────────────────────────────────────────────────────
def _rep_para(para, rep: dict):
    """Replace text in each run individually — preserves bold/italic/colour per run.
    Previously this crushed all runs into run[0] which made everything bold.
    Now each run is updated in-place so formatting is never disturbed."""
    for run in para.runs:
        t = run.text
        for k, v in rep.items():
            if k and v is not None and k in t:
                t = t.replace(k, v)
        if t != run.text:
            run.text = t


def _rep_shape(shape, rep: dict):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            _rep_para(p, rep)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    _rep_para(p, rep)
    if shape.shape_type == 6:
        for s in shape.shapes:
            _rep_shape(s, rep)


def rep_slide(slide, rep: dict):
    for shape in slide.shapes:
        _rep_shape(shape, rep)


def clean_apostrophes(prs):
    """
    Fix all apostrophe-space gaps across every text run in the presentation.

    Two root causes:
    1. Template typos: e.g. 'EIIL\u2019 privacy' (curly-apos + space, 's' missing)
       → after name replace: 'SGD Pharma\u2019 privacy' — the space survives
    2. Run-split: run[j] ends with company name (trailing space), run[j+1] starts
       with apostrophe → renders as 'SGD Pharma 's'

    Fix strategy:
    a) Within each run: collapse 'word \u2019s' → 'word\u2019s' and 'word \'s' → 'word\'s'
    b) Across run boundaries: if run[j] ends with a letter/space and run[j+1] starts
       with an apostrophe, strip trailing space from run[j]
    """
    APOS = ("\u2019", "\u2018", "'")   # curly-right, curly-left, straight

    def fix_run_text(text):
        for ap in APOS:
            # Pattern: word-char + space + apostrophe → word-char + apostrophe (no space)
            text = re.sub(r'(\w) ' + re.escape(ap), r'\1' + ap, text)
        return text

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = list(para.runs)
                # (a) Fix within each run
                for run in runs:
                    fixed = fix_run_text(run.text)
                    if fixed != run.text:
                        run.text = fixed
                # (b) Fix across run boundaries
                for j in range(len(runs) - 1):
                    for ap in APOS:
                        if runs[j + 1].text.startswith(ap):
                            # Strip trailing space from run[j]
                            if runs[j].text.endswith(" "):
                                runs[j].text = runs[j].text.rstrip(" ")


def set_para_text(slide, shape_name: str, fragment: str, new_text: str) -> bool:
    """Find paragraph whose combined run-text contains fragment; replace with new_text."""
    for shape in slide.shapes:
        if shape_name and shape.name != shape_name:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if fragment in full:
                if para.runs:
                    para.runs[0].text = new_text
                    for r in para.runs[1:]:
                        r.text = ""
                return True
    return False


# ─────────────────────────────────────────────────────────────
# FIX SLIDE 4 BULLET CONSISTENCY
# ─────────────────────────────────────────────────────────────
def fix_slide4_bullets(slide):
    """
    Ensure TextBox 12 right-column bullet paragraphs render as Arial • circles.

    CRITICAL OOXML ORDER inside <a:pPr>:
      spcBef → spcAft → buClr → buSzPct → buFont → buChar → tabLst → defRPr → extLst
    Appending to the end puts bullets AFTER defRPr → PowerPoint ignores them → no dots.
    We must INSERT before defRPr using the element's index in pPr's child list.
    """
    BULLET_TAGS = ("buNone", "buClrTx", "buClr", "buSzTx", "buSzPct",
                   "buSzClamp", "buFont", "buFontTx", "buChar", "buAutoNum")

    for shape in slide.shapes:
        if shape.name != "TextBox 12":
            continue
        # Snapshot paragraphs to a plain list — avoids lxml proxy identity issues
        paras = list(shape.text_frame.paragraphs)
        for i, para in enumerate(paras):
            if i < 2:                               # P0=heading, P1=intro — no bullet
                continue
            pPr = para._p.find(f"{{{ANS}}}pPr")
            if pPr is None:
                continue

            # ── Remove ALL existing bullet-related elements ──────
            for tag in BULLET_TAGS:
                el = pPr.find(f"{{{ANS}}}{tag}")
                if el is not None:
                    pPr.remove(el)

            # ── Build the 4 elements matching the left column ───
            bu_clr  = etree.fromstring(
                f'<a:buClr xmlns:a="{ANS}"><a:srgbClr val="3C3D3E"/></a:buClr>'
            )
            bu_sz   = etree.fromstring(
                f'<a:buSzPct xmlns:a="{ANS}" val="100000"/>'
            )
            bu_font = etree.fromstring(
                f'<a:buFont xmlns:a="{ANS}" typeface="Arial" '
                f'panose="020B0604020202020204" pitchFamily="34" charset="0"/>'
            )
            bu_char = etree.fromstring(
                f'<a:buChar xmlns:a="{ANS}" char="\u2022"/>'
            )

            # ── INSERT before <a:defRPr> so OOXML order is valid ─
            pPr_children = list(pPr)               # snapshot children list
            defRPr = pPr.find(f"{{{ANS}}}defRPr")
            if defRPr is not None:
                ins_idx = pPr_children.index(defRPr)
                # Insert in reverse order so they land in the correct sequence
                for el in (bu_char, bu_font, bu_sz, bu_clr):
                    pPr.insert(ins_idx, el)
            else:
                for el in (bu_clr, bu_sz, bu_font, bu_char):
                    pPr.append(el)

            # ── Ensure hanging-indent attributes are present ─────
            if not pPr.get("marL"):   pPr.set("marL",   "171450")
            if not pPr.get("indent"): pPr.set("indent", "-171450")


# ─────────────────────────────────────────────────────────────
# MAIN BUILD
# ─────────────────────────────────────────────────────────────
def build_presentation(pptx_bytes: bytes, company_name: str,
                       company_short: str, ai: dict) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))

    # ── Global name replacements on every slide ───────────────
    gmap = {
        "Eveready Industries India Ltd. (EIIL)": company_name,
        "Eveready Industries India Ltd":         company_name.split("(")[0].strip().rstrip(","),
        "Eveready Industries":                   company_name.split("(")[0].strip().rstrip(","),
        # Straight apostrophe variants
        " EIIL'":   f" {company_short}'",
        " EIIL's":  f" {company_short}'s",
        # Curly-apostrophe variants (U+2019) — used throughout the template
        " EIIL\u2019s": f" {company_short}\u2019s",
        " EIIL\u2019":  f" {company_short}\u2019s",   # template typo: missing 's' after apostrophe
        "EIIL\u2019s":  f"{company_short}\u2019s",
        "EIIL\u2019 ":  f"{company_short}\u2019s ",   # template typo in same run (no leading space)
        "EIIL\u2019":   f"{company_short}\u2019s",    # bare curly apos without trailing space
        "(EIIL)":        f"({company_short})",
        "EIIL":          company_short,
    }
    for slide in prs.slides:
        rep_slide(slide, gmap)

    # ── Slide 4 (index 3) ────────────────────────────────────
    if len(prs.slides) > 3:
        s4 = prs.slides[3]

        # ── TextBox 8: Top company description ───────────────
        # Run structure: R0(bold)=company name, R1(normal)=body
        # Global replace already updated R0 (EIIL→company_name).
        # Write AI body to R1 ONLY — never touch R0 (preserves bold on name only).
        if ai.get("s4_desc"):
            for shape in s4.shapes:
                if shape.name != "TextBox 8": continue
                runs = shape.text_frame.paragraphs[0].runs
                if len(runs) >= 2:
                    runs[1].text = ai["s4_desc"]
                    for r in runs[2:]: r.text = ""
                break

        # ── TextBox 3 P0: Scope paragraph ────────────────────
        # Run structure: R0(bold)=company short, R1(normal)=fixed DPDPA text,
        #                R2(bold)=fixed "Digital Personal Data Protection Act, 2023",
        #                R3(normal)=" and applicable Rules...landscape across [OPS]."
        # Global replace already updated R0 (EIIL→company_short).
        # ONLY update R3 with AI operations phrase — R1 and R2 stay exactly as-is.
        if ai.get("s4_ops"):
            for shape in s4.shapes:
                if shape.name != "TextBox 3": continue
                runs = shape.text_frame.paragraphs[0].runs
                if len(runs) >= 4:
                    ops_text = ai["s4_ops"].strip()
                    # Ensure it starts with " and applicable Rules..."
                    if not ops_text.lower().startswith("and "):
                        ops_text = "and applicable Rules, calibrated to its people, process and technology landscape across " + ops_text
                    # Prepend space (R3 starts with a space in original)
                    if not ops_text.startswith(" "):
                        ops_text = " " + ops_text
                    # Ensure ends with period
                    if not ops_text.rstrip().endswith("."):
                        ops_text = ops_text.rstrip() + "."
                    runs[3].text = ops_text
                    for r in runs[4:]: r.text = ""
                break

        # ── TextBox 3 P2–P8: Scope bullets ───────────────────
        bullet_frags = [
            "Conduct an enterprise",
            "Assess privacy, information security",
            "Evaluate governance",
            "Design and operationalize",
            "Support rollout",
            "Coordinate remediation",
            "Deliver role",
        ]
        for frag, new_b in zip(bullet_frags, ai.get("s4_bullets", [])):
            if new_b:
                set_para_text(s4, "TextBox 3", frag, new_b)

        # ── TextBox 12: Fix bullets then write AI text ────────
        # P1 (intro) is NEVER rewritten — it's generic and already perfect:
        # "We help embed a trust-first approach...DPDPA and it's Rules."
        # P0=heading (skip), P1=intro (skip), P2–P7=6 bullets (AI)
        fix_slide4_bullets(s4)

        s4r = ai.get("s4_right", {})
        if s4r:
            for shape in s4.shapes:
                if shape.name != "TextBox 12": continue
                paras = list(shape.text_frame.paragraphs)
                # P2..P7 — 6 bullets only (P0 and P1 are never touched)
                for i, key in enumerate(["b1","b2","b3","b4","b5","b6"], start=2):
                    if s4r.get(key) and len(paras) > i:
                        p = paras[i]
                        if p.runs:
                            p.runs[0].text = s4r[key]
                            for r in p.runs[1:]: r.text = ""
                break

    # ── Slide 11 (index 10) ──────────────────────────────────
    if len(prs.slides) > 10:
        set_para_text(prs.slides[10], "Rectangle 1",
                      "For this engagement", ai["s11"])

    # ── Slide 17 (index 16) ──────────────────────────────────
    if len(prs.slides) > 16:
        s17 = prs.slides[16]
        lc  = ai.get("s17_lifecycle", {})
        for shp_name, key in [
            ("Rectangle 41", "collection"),
            ("Rectangle 8",  "use_processing"),
            ("Rectangle 9",  "storage"),
            ("Rectangle 10", "sharing"),
            ("Rectangle 11", "retention"),
            ("Rectangle 12", "disposal"),
        ]:
            if lc.get(key):
                set_para_text(s17, shp_name, "We will", lc[key])

    # ── Slides 12, 14, 19 – targeted sentence replacement ────
    s12_b1     = ai.get("s12_b1", "")
    s12_b2     = ai.get("s12_b2", "")
    s14_b1     = ai.get("s14_b1", "")
    s14_b2     = ai.get("s14_b2", "")
    s19_notice = ai.get("s19_notice", "")

    # Fragment anchors (U+2011 = non-breaking hyphen used in template)
    FRAG_B1      = "enterprise\u2011wide privacy applicability"
    FRAG_B1b     = "enterprise-wide privacy applicability"   # fallback
    FRAG_B2      = "information security and regulatory risks"
    FRAG_NOTICE  = "Privacy Notice and Consent Notice"

    def replace_in_slide(slide, fragment, new_text):
        """Replace paragraph containing fragment in any shape. Returns True if found."""
        if not new_text:
            return False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if fragment in full:
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]:
                            r.text = ""
                    return True
        return False

    # Slide 12 (index 11) — narrow 264pt box, 22w limit
    if len(prs.slides) > 11:
        sl12 = prs.slides[11]
        if s12_b1:
            if not replace_in_slide(sl12, FRAG_B1, s12_b1):
                replace_in_slide(sl12, FRAG_B1b, s12_b1)
        if s12_b2:
            replace_in_slide(sl12, FRAG_B2, s12_b2)

    # Slide 14 (index 13) — wide 661pt box, 28w limit
    if len(prs.slides) > 13:
        sl14 = prs.slides[13]
        if s14_b1:
            if not replace_in_slide(sl14, FRAG_B1, s14_b1):
                replace_in_slide(sl14, FRAG_B1b, s14_b1)
        if s14_b2:
            replace_in_slide(sl14, FRAG_B2, s14_b2)

    # Slide 19 (index 18) — Privacy Notice sentence
    if len(prs.slides) > 18 and s19_notice:
        replace_in_slide(prs.slides[18], FRAG_NOTICE, s19_notice)

    # ── Final pass: remove all apostrophe-space gaps in every slide ─
    clean_apostrophes(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────
# GENERATE FLOW
# ─────────────────────────────────────────────────────────────
if generate_btn:
    # ── Validation ───────────────────────────────────────────
    errs = []
    if not groq_api_key:    errs.append("🔑 Groq API Key required.")
    if not company_name:    errs.append("🏢 Full company name required.")
    if not company_short:   errs.append("🏷️ Abbreviation required.")
    if not company_website: errs.append("🌐 Website URL required.")
    if not uploaded_ppt:    errs.append("📁 Template PPTX required.")
    for e in errs:
        st.error(e)
    if errs:
        st.stop()

    pptx_bytes = uploaded_ppt.read()
    client     = Groq(api_key=groq_api_key)

    # ── Step 1: Website scrape ────────────────────────────────
    with st.status("🌐 Scraping company website…", expanded=False) as sts:
        web = scrape_website(company_website)
        ok  = not web.startswith("[Scrape")
        sts.update(
            label="✅ Website scraped" if ok
                  else "⚠️ Could not scrape website — proceeding with company name only",
            state="complete" if ok else "error",
        )

    # ── Step 2: AI content generation ────────────────────────
    ai = {}
    def safe(key, prompt, max_tok=1000, fallback=""):
        try:
            ai[key] = groq_call(client, prompt, max_tok)
        except Exception as e:
            ai[key] = fallback
            st.warning(f"{key}: {e}")

    with st.status("🤖 Generating AI content for slides…", expanded=True) as sts:

        # ── Pre-step 1: Extract full company profile ──────────
        st.write("🔍 Extracting company profile from website…")
        profile = {}
        try:
            raw_profile = groq_call(client,
                                    P_COMPANY_PROFILE.format(
                                        company_name=company_name,
                                        company_short=company_short,
                                        website_text=web[:3000]),
                                    max_tokens=600)
            raw_profile = re.sub(r"^```(?:json)?", "", raw_profile).strip()
            raw_profile = re.sub(r"```$", "", raw_profile).strip()
            profile = json.loads(raw_profile)
        except Exception as e:
            st.warning(f"Profile extraction: {e}")

        # Safe accessors with sensible fallbacks
        industry      = profile.get("industry",             "business services")
        biz_model     = profile.get("business_model",       f"Predominantly B2B, serving enterprise clients through direct channels")
        svc_lines     = profile.get("service_lines",        company_short + "'s core service lines")
        key_sectors   = profile.get("key_sectors",          "enterprise clients across multiple sectors")
        key_funcs     = profile.get("key_functions",        "Operations, Finance, HR, IT, Compliance, Client Delivery")
        key_systems   = profile.get("key_systems",          "ERP, CRM, cloud platforms and analytics tools")
        data_types    = profile.get("data_types",           "employee data, client data, operational data")
        partner_types = profile.get("partner_types",        "clients, vendors, sub-contractors, regulators")
        geo_footprint = profile.get("geographic_footprint", "pan-India with global operations")
        channel_desc  = profile.get("channel_description",  "direct B2B sales and partner channels")
        ai["biz_model"] = biz_model
        ai["profile"]   = profile

        st.write("📝 Slide 4 — Company description body…")
        safe("s4_desc",
             P_DESC.format(company_name=company_name,
                           company_short=company_short,
                           industry=industry,
                           business_model=biz_model,
                           service_lines=svc_lines,
                           key_sectors=key_sectors,
                           key_systems=key_systems,
                           geographic_footprint=geo_footprint,
                           channel_description=channel_desc),
             max_tok=500,
             fallback=(
                 f"is a leading {industry} company, operating through a diversified model "
                 f"spanning {svc_lines} across domestic and select international markets. "
                 f"{biz_model.rstrip('.')}. "
                 f"{company_short} enables end-to-end service delivery through technology-driven "
                 f"quality systems and integrated operations, ensuring reliable and compliant "
                 f"delivery across diverse client segments."
             ))
        if ai.get("s4_desc"):
            body = ai["s4_desc"].strip()
            for prefix in [company_name, company_short]:
                if body.lower().startswith(prefix.lower()):
                    body = body[len(prefix):].lstrip(" ,")
            if not body.lower().startswith("is "):
                body = "is " + body
            ai["s4_desc"] = " ".join(body.split()[:124])

        st.write("📝 Slide 4 — Scope operations phrase…")
        safe("s4_ops",
             P_SCOPE_OPS.format(company_name=company_name,
                                company_short=company_short,
                                industry=industry,
                                service_lines=svc_lines,
                                key_sectors=key_sectors),
             max_tok=80,
             fallback=f"and applicable Rules, calibrated to its people, process and technology landscape across {company_short}'s {industry.lower()} operations.")

        st.write("📝 Slide 4 — Scope bullets (7)…")
        try:
            raw_b = groq_call(client,
                              P_BULLETS.format(company_name=company_name,
                                               company_short=company_short,
                                               industry=industry,
                                               service_lines=svc_lines,
                                               key_sectors=key_sectors,
                                               key_functions=key_funcs,
                                               key_systems=key_systems,
                                               partner_types=partner_types),
                              max_tokens=1400)
            lines = [l.strip() for l in raw_b.split("\n") if l.strip()]
            lines = [re.sub(r"^[\d]+[.)]\s*", "", l).lstrip("•–-").strip() for l in lines]
            ai["s4_bullets"] = lines[:7]
        except Exception as e:
            ai["s4_bullets"] = []
            st.warning(f"s4_bullets: {e}")

        st.write("📝 Slide 4 — Right-side 'How We Will Help' bullets…")
        S4R_LIMITS = {"b1":28, "b2":22, "b3":16, "b4":25, "b5":25, "b6":30}
        try:
            raw_r = groq_call(client,
                              P_S4_RIGHT.format(company_name=company_name,
                                                company_short=company_short,
                                                industry=industry,
                                                business_model=biz_model,
                                                service_lines=svc_lines,
                                                key_sectors=key_sectors),
                              max_tokens=800)
            raw_r = re.sub(r"^```(?:json)?", "", raw_r).strip()
            raw_r = re.sub(r"```$",          "", raw_r).strip()
            s4r = json.loads(raw_r)
            for k, lim in S4R_LIMITS.items():
                if k in s4r:
                    s4r[k] = " ".join(s4r[k].split()[:lim])
            ai["s4_right"] = s4r
        except Exception as e:
            ai["s4_right"] = {}
            st.warning(f"s4_right: {e}")

        st.write("📝 Slide 11 — Operating model paragraph…")
        safe("s11",
             P_S11.format(company_name=company_name,
                          company_short=company_short,
                          industry=industry,
                          business_model=biz_model,
                          service_lines=svc_lines,
                          key_functions=key_funcs,
                          geographic_footprint=geo_footprint,
                          channel_description=channel_desc),
             max_tok=250,
             fallback=(
                 f"For this engagement, the privacy compliance model will be applied exclusively "
                 f"to the internal functions, processes and governance structures of {company_name}, "
                 f"supporting its {svc_lines}, which {biz_model.rstrip('.')}, with limited B2C "
                 f"personal data processing through digital platform and client service interactions. "
                 f"This approach ensures a focused effort on strengthening {company_short}'s internal "
                 f"privacy governance and compliance capabilities, aligned with applicable regulatory "
                 f"requirements, its operating model and its {geo_footprint} operational footprint."
             ))
        if ai.get("s11"):
            words = ai["s11"].split()
            if len(words) > 82:
                # Cut at last full stop within 82 words, never mid-sentence
                candidate = " ".join(words[:82])
                last_dot = candidate.rfind('.')
                ai["s11"] = candidate[:last_dot + 1].strip() if last_dot > 0 else candidate + '.'

        st.write("📝 Slide 17 — Data Lifecycle (6 sections)…")
        S17_LIMITS = {
            "collection": 61, "use_processing": 64, "storage": 48,
            "sharing": 36, "retention": 43, "disposal": 47,
        }
        try:
            raw17 = groq_call(client,
                              P_S17.format(company_name=company_name,
                                           company_short=company_short,
                                           industry=industry,
                                           service_lines=svc_lines,
                                           key_sectors=key_sectors,
                                           key_functions=key_funcs,
                                           key_systems=key_systems,
                                           data_types=data_types,
                                           partner_types=partner_types,
                                           geographic_footprint=geo_footprint),
                              max_tokens=3000)
            raw17 = re.sub(r"^```(?:json)?", "", raw17).strip()
            raw17 = re.sub(r"```$", "", raw17).strip()
            ai["s17_lifecycle"] = json.loads(raw17)
            for k, limit in S17_LIMITS.items():
                if k in ai["s17_lifecycle"]:
                    words = ai["s17_lifecycle"][k].split()
                    if len(words) > limit:
                        ai["s17_lifecycle"][k] = " ".join(words[:limit])
        except Exception as e:
            ai["s17_lifecycle"] = {}
            st.warning(f"s17_lifecycle: {e}")

        st.write("📝 Slides 12, 14, 19 — Tailoring operational sentences…")

        def smart_trim(text: str, max_words: int) -> str:
            """
            Trim text to max_words but ALWAYS end at a complete sentence or clause.
            Never cuts mid-sentence. Priority order:
            1. If text fits within max_words → return as-is (ensure ends with '.')
            2. Find the last full stop (.) before the word limit → cut there
            3. Find the last comma before the word limit → cut there, add '.'
            4. Hard cut at max_words, add '.'
            """
            words = text.split()
            if len(words) <= max_words:
                t = " ".join(words)
                return t if t.endswith('.') else t.rstrip(',') + '.'

            # Rebuild within limit
            candidate = " ".join(words[:max_words])

            # Try to find last full stop within candidate
            last_dot = candidate.rfind('.')
            if last_dot > len(candidate) * 0.5:   # dot is at least halfway through
                return candidate[:last_dot + 1].strip()

            # Try last comma (end of a clause)
            last_comma = candidate.rfind(',')
            if last_comma > len(candidate) * 0.4:
                return candidate[:last_comma].strip() + '.'

            # Hard cut — at least add a period
            return candidate.rstrip(',').rstrip() + '.'

        # Slide 12: narrow box (264pt, 9pt font) — 32 word capacity
        try:
            raw_s12 = groq_call(client,
                                P_S12_BULLETS.format(company_name=company_name,
                                                     company_short=company_short,
                                                     industry=industry,
                                                     key_functions=key_funcs,
                                                     key_systems=key_systems,
                                                     partner_types=partner_types,
                                                     service_lines=svc_lines),
                                max_tokens=250)
            s12_lines = [re.sub(r"^[\d]+[.)]\s*","",l).lstrip("•–-").strip()
                         for l in raw_s12.strip().split("\n") if l.strip()]
            ai["s12_b1"] = smart_trim(s12_lines[0], 32) if len(s12_lines)>0 else ""
            ai["s12_b2"] = smart_trim(s12_lines[1], 32) if len(s12_lines)>1 else ""
        except Exception as e:
            ai["s12_b1"] = ai["s12_b2"] = ""
            st.warning(f"s12_bullets: {e}")

        # Slide 14: wide box (661pt, 11pt font) — hard limit 28 words per sentence
        try:
            raw_s14 = groq_call(client,
                                P_S14_BULLETS.format(company_name=company_name,
                                                     company_short=company_short,
                                                     industry=industry,
                                                     key_functions=key_funcs,
                                                     key_systems=key_systems,
                                                     partner_types=partner_types),
                                max_tokens=250)
            s14_lines = [re.sub(r"^[\d]+[.)]\s*","",l).lstrip("•–-").strip()
                         for l in raw_s14.strip().split("\n") if l.strip()]
            ai["s14_b1"] = smart_trim(s14_lines[0], 30) if len(s14_lines)>0 else ""
            ai["s14_b2"] = smart_trim(s14_lines[1], 30) if len(s14_lines)>1 else ""
        except Exception as e:
            ai["s14_b1"] = ai["s14_b2"] = ""
            st.warning(f"s14_bullets: {e}")

        safe("s19_notice",
             P_S19_NOTICE.format(company_name=company_name,
                                 company_short=company_short,
                                 industry=industry,
                                 key_functions=key_funcs,
                                 key_systems=key_systems,
                                 service_lines=svc_lines),
             max_tok=120,
             fallback=f"Prepare tailored Privacy Notice and Consent Notice for {company_short}'s touchpoints covering operations, quality, HR, finance, IT, enterprise platforms and cloud/SaaS applications.")

        sts.update(label="✅ All AI content generated", state="complete")

    # ── Step 3: Build PPTX ───────────────────────────────────
    with st.status("📝 Applying changes to PPTX…", expanded=False) as sts:
        try:
            output = build_presentation(pptx_bytes, company_name, company_short, ai)
            sts.update(label="✅ PPTX ready!", state="complete")
        except Exception as e:
            sts.update(label="❌ Failed", state="error")
            st.error(str(e))
            import traceback; st.code(traceback.format_exc())
            st.stop()

    st.success("🎉 Proposal generated successfully!")

    # ── Preview ──────────────────────────────────────────────
    with st.expander("🔍 Preview AI-generated content"):
        if ai.get("profile"):
            p = ai["profile"]
            st.markdown("**🏢 Extracted Company Profile:**")
            st.markdown(f"""
| Field | Value |
|---|---|
| **Industry** | {p.get('industry','')} |
| **Business Model** | {p.get('business_model','')} |
| **Service Lines** | {p.get('service_lines','')} |
| **Key Sectors** | {p.get('key_sectors','')} |
| **Key Functions** | {p.get('key_functions','')} |
| **Key Systems** | {p.get('key_systems','')} |
| **Data Types** | {p.get('data_types','')} |
| **Partner Types** | {p.get('partner_types','')} |
| **Geography** | {p.get('geographic_footprint','')} |
| **Channel** | {p.get('channel_description','')} |
""")
        st.markdown("**Slide 4 – Company Description:**")
        st.info(ai.get("s4_desc", ""))
        st.markdown("**Slide 4 – Scope Operations Phrase:**")
        st.info(ai.get("s4_ops", ""))
        if ai.get("s4_bullets"):
            st.markdown("**Slide 4 – Scope Bullets (Left):**")
            for b in ai["s4_bullets"]:
                st.write(f"• {b}")
        if ai.get("s4_right"):
            st.markdown("**Slide 4 – How We Will Help (Right):**")
            r = ai["s4_right"]
            for k in ["b1","b2","b3","b4","b5","b6"]:
                if r.get(k): st.write(f"• {r[k]}")
        st.markdown("**Slide 11 – Operating Model Paragraph:**")
        st.info(ai.get("s11", ""))
        if ai.get("s17_lifecycle"):
            st.markdown("**Slide 17 – Data Lifecycle:**")
            for k, v in ai["s17_lifecycle"].items():
                st.write(f"**{k.replace('_', ' ').title()}:** {v}")

    # ── Download ─────────────────────────────────────────────
    safe_name = re.sub(r"[^\w\s-]", "", company_name).strip().replace(" ", "_")[:40]
    filename  = f"Proposal_Data_Privacy_{safe_name}_March2026.pptx"
    st.download_button(
        label="⬇️ Download Personalised Proposal",
        data=output,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        type="primary",
    )
    st.caption(
        "ℹ️ Slide 5 (Scope of Review) is intentionally left unchanged — "
        "please fill in company-specific details manually."
    )

# ── Landing state ─────────────────────────────────────────────
else:
    st.info("👈 Fill in all details in the sidebar and upload the template, then click **Generate Proposal**.")
    with st.expander("📖 What changes per slide"):
        st.markdown("""
| Slide | Change | Method |
|---|---|---|
| **1** | Company name in title | Auto-replace |
| **4** | Company description (detailed 3-sentence paragraph) + scope paragraph + 7 bullets | Groq AI |
| **4** | Right-side bullet style fixed to ● circle (matching left side) | Auto-fix |
| **5** | **Not changed** — fill manually | Human |
| **11** | Privacy Operating Model paragraph | Groq AI |
| **12, 14, 19** | All EIIL/Eveready name references | Auto-replace |
| **17** | All 6 Data Lifecycle sections | Groq AI |

All visuals, colours, fonts, charts, team slides and Protiviti content remain untouched.
""")
